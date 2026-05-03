import os
import tempfile
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List, Dict, Any

from rag_pipeline import RAGPipeline
from llm_provider import LLMProvider

app = FastAPI(title="Naive RAG Pipeline API")

# Add CORS middleware to allow React frontend to communicate with backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

pipeline = RAGPipeline(vectorstore_dir="../vectorstore")
documents_text = ""

class QueryRequest(BaseModel):
    query: str
    provider: str = "groq"
    model: str = "llama-3.1-8b-instant"
    api_key: Optional[str] = None
    k: int = 4

@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    global documents_text
    
    valid_extensions = ('.txt', '.pdf', '.docx', '.xlsx', '.pptx', '.csv')
    if not file.filename.endswith(valid_extensions):
        raise HTTPException(status_code=400, detail=f"Only {', '.join(valid_extensions)} files are supported.")
        
    try:
        content = ""
        # Read the file
        file_bytes = await file.read()
        
        if file.filename.endswith('.pdf'):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
                
            with pdfplumber.open(temp_file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
            os.unlink(temp_file_path)
            
        elif file.filename.endswith('.docx'):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            doc = docx.Document(temp_file_path)
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            os.unlink(temp_file_path)
            
        elif file.filename.endswith('.xlsx'):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            wb = openpyxl.load_workbook(temp_file_path, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        content += row_text + "\n"
            os.unlink(temp_file_path)
            
        elif file.filename.endswith('.pptx'):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            prs = Presentation(temp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            os.unlink(temp_file_path)
            
        elif file.filename.endswith('.csv'):
            import csv
            with tempfile.NamedTemporaryFile(delete=False, suffix=".csv") as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            with open(temp_file_path, mode='r', encoding='utf-8', errors='ignore') as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    row_text = " ".join([str(cell) for cell in row])
                    if row_text.strip():
                        content += row_text + "\n"
            os.unlink(temp_file_path)
            
        else:
            content = file_bytes.decode('utf-8')
            
        if not content.strip():
            raise HTTPException(status_code=400, detail="Document is empty or text could not be extracted.")
            
        documents_text = content
        
        return {
            "success": True,
            "filename": file.filename,
            "text_preview": content[:500] + "..." if len(content) > 500 else content,
            "char_count": len(content)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process_pipeline")
async def process_pipeline():
    global documents_text
    if not documents_text:
        raise HTTPException(status_code=400, detail="No document uploaded yet.")
        
    try:
        # 1. Chunking
        chunk_res = pipeline.chunk_text(documents_text)
        
        # 2. Embeddings
        embed_res = pipeline.create_embeddings(chunk_res["chunks"])
        
        return {
            "success": True,
            "chunking": chunk_res,
            "embeddings": embed_res
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query")
async def run_query(request: QueryRequest):
    try:
        # 1. Retrieval
        retrieval_res = pipeline.retrieve(request.query, k=request.k)
        
        # 2. Prompt Construction
        context_texts = [res["content"] for res in retrieval_res["results"]]
        context = "\n\n---\n\n".join(context_texts)
        
        prompt = f"""You are a helpful assistant. Use the following pieces of retrieved context to answer the question. 
If you don't know the answer, just say that you don't know. 

Context:
{context}

Question:
{request.query}

Answer:"""
        
        # 3. LLM Generation
        llm = LLMProvider(provider=request.provider, api_key=request.api_key, model=request.model)
        generation_res = llm.generate(prompt)
        
        if not generation_res["success"]:
            raise HTTPException(status_code=500, detail=generation_res.get("error", "LLM generation failed"))
            
        return {
            "success": True,
            "retrieval": retrieval_res,
            "prompt": prompt,
            "generation": generation_res
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
